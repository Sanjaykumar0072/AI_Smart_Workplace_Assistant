import fitz  # pymupdf
from docx import Document
from pptx import Presentation
from io import BytesIO


def read_file(uploaded_file):
    file_name = uploaded_file.name.lower()
    file_bytes = uploaded_file.read()
    file_stream = BytesIO(file_bytes)

    try:
        # -------------------------
        # PDF
        # -------------------------
        if file_name.endswith(".pdf"):
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            return text if text else "No readable text found in PDF."

        # -------------------------
        # Word (.docx)
        # -------------------------
        elif file_name.endswith(".docx"):
            doc = Document(file_stream)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text if text else "No readable text found in Word document."

        # -------------------------
        # PowerPoint (.pptx)
        # -------------------------
        elif file_name.endswith(".pptx"):
            prs = Presentation(file_stream)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text if text else "No readable text found in PowerPoint."

        # -------------------------
        # Text File
        # -------------------------
        elif file_name.endswith(".txt"):
            return file_bytes.decode("utf-8")

        # -------------------------
        # Unsupported
        # -------------------------
        else:
            return "Unsupported file format."

    except Exception as e:
        return f"Error reading file: {str(e)}"